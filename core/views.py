import re
from django.shortcuts import render, redirect, get_object_or_404
import zipfile
import io
from django.http import HttpResponse
from django.contrib.auth import login, logout, authenticate
from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth.forms import AuthenticationForm
from django.contrib.auth.models import User
from django.contrib.auth.decorators import login_required
from django.contrib import messages
from .models import ServiceOrder, OrderFile, Profile, SiteSetting, ContactMessage, OrderChat, Review, CaseStudy, AgencyStat, TeamMember, Freelancer, FreelancerChat, FreelancerNotification
from django.conf import settings
from django.core.mail import send_mail
from django.utils import timezone
import datetime
import random
import logging
# Re-import firebase_admin for Google Auth
import firebase_admin
from firebase_admin import auth as firebase_auth
import fitz  # PyMuPDF
import tempfile
import os
from pptx import Presentation
import pdfplumber
import pandas as pd
import openpyxl
from pptx.util import Inches

logger = logging.getLogger(__name__)

# --- HELPER FUNCTIONS ---
def get_user_by_email_or_phone(identifier):
    try:
        user = User.objects.get(email=identifier)
        return user
    except User.DoesNotExist:
        pass
    
    try:
        profile = Profile.objects.get(phone_number=identifier)
        return profile.user
    except Profile.DoesNotExist:
        pass
    return None

def verify_google_token(id_token):
    try:
        decoded_token = firebase_auth.verify_id_token(id_token)
        return decoded_token
    except Exception as e:
        logger.error(f"Error verifying Firebase token: {e}")
        return None

def validate_input(email, phone):
    # Regex for Email
    email_regex = r'^[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,}$'
    if not re.match(email_regex, email):
        return False, "Invalid email address format."
    
    # Regex for Phone (Allow +91 or just 10 digits)
    phone_regex = r'^(\+91[\-\s]?)?[0]?(91)?[6789]\d{9}$'
    if not re.match(phone_regex, phone):
        return False, "Invalid phone number. Must be a valid Indian mobile number."
        
    return True, ""

# --- VIEWS ---

def home(request):
    if request.user.is_authenticated:
        return redirect('dashboard')
    reviews = Review.objects.all().order_by('-created_at')
    return render(request, 'core/home.html', {'reviews': reviews})

def tools_list(request):
    return render(request, 'core/tools_list.html')

def signup_view(request):
    if request.method == 'POST':
        # Check for Google Sign-In first
        id_token = request.POST.get('google_id_token')
        if id_token:
            decoded_token = verify_google_token(id_token)
            if decoded_token:
                email = decoded_token.get('email')
                name = decoded_token.get('name')
                # Google Users verified by default
                
                # Check if user exists
                user = User.objects.filter(email=email).first()
                if not user:
                    # Create new Google User
                    username = email.split('@')[0] + str(random.randint(1000,9999))
                    user = User.objects.create_user(username=username, email=email)
                    user.first_name = name
                    user.save()
                    Profile.objects.create(user=user) # Empty phone for google initially
                
                login(request, user)
                return redirect('dashboard')
            else:
                 messages.error(request, "Google Sign-In failed.")
                 return redirect('signup')

        # Standard Signup
        full_name = request.POST.get('full_name')
        email = request.POST.get('email')
        phone = request.POST.get('phone')
        password = request.POST.get('password')
        confirm_password = request.POST.get('confirm_password')

        # 1. Regex Validation
        is_valid, error_msg = validate_input(email, phone)
        if not is_valid:
            messages.error(request, error_msg)
            return redirect('signup')

        if password != confirm_password:
            messages.error(request, "Passwords do not match!")
            return redirect('signup')
        
        if User.objects.filter(email=email).exists():
            messages.error(request, "Email already registered!")
            return redirect('signup')
            
        if Profile.objects.filter(phone_number=phone).exists():
            messages.error(request, "Phone number already registered!")
            return redirect('signup')
        
        # Create User
        try:
            # We use phone number as username to ensure uniqueness
            user = User.objects.create_user(
                username=phone, 
                email=email,
                password=password
            )
            user.first_name = full_name
            user.save()
            
            Profile.objects.create(user=user, phone_number=phone)
            
            login(request, user)
            messages.success(request, "Account created successfully!")
            return redirect('dashboard')
            
        except Exception as e:
            messages.error(request, f"Error creating account: {e}")
            return redirect('signup')

    return render(request, 'core/signup.html')

def verify_otp(request):
    # Deprecated view, redirecting to home
    return redirect('home')

def login_view(request):
    if request.method == 'POST':
        # Check for Google Sign-In first
        id_token = request.POST.get('google_id_token')
        if id_token:
            decoded_token = verify_google_token(id_token)
            if decoded_token:
                email = decoded_token.get('email')
                
                # Check if user exists
                user = User.objects.filter(email=email).first()
                if user:
                    login(request, user)
                    return redirect('dashboard')
                else:
                    # Optional: Auto-create account on login if it doesn't exist?
                    # For now, let's redirect to signup with a message or just create it.
                    # User requested "Sign in with Google" which usually implies auto-signup.
                    name = decoded_token.get('name')
                    username = email.split('@')[0] + str(random.randint(1000,9999))
                    user = User.objects.create_user(username=username, email=email)
                    user.first_name = name
                    user.save()
                    Profile.objects.create(user=user)
                    login(request, user)
                    return redirect('dashboard')
            else:
                 messages.error(request, "Google Sign-In failed.")
                 return redirect('login')

        identifier = request.POST.get('identifier')
        password = request.POST.get('password')
        
        user = get_user_by_email_or_phone(identifier)
        
        if user:
            user = authenticate(request, username=user.username, password=password)
            if user is not None:
                login(request, user)
                return redirect('dashboard')
        
        messages.error(request, "Invalid Email/Phone or Password")
    
    return render(request, 'core/login.html')

def logout_view(request):
    logout(request)
    return redirect('home')

@login_required
def dashboard(request):
    orders = ServiceOrder.objects.filter(user=request.user).order_by('-created_at')
    
    stats = {
        'total': orders.count(),
        'pending': orders.filter(status='pending').count(),
        'in_progress': orders.filter(status='in_progress').count(),
        'completed': orders.filter(status='completed').count(),
    }
    
    return render(request, 'core/dashboard.html', {'orders': orders, 'stats': stats})

@login_required
def create_order(request):
    if request.method == 'POST':
        service_type = request.POST.get('service_type')
        title = request.POST.get('title')
        description = request.POST.get('description')
        phone_number = request.POST.get('phone_number')
        request_call = request.POST.get('request_call') == 'on'
        files = request.FILES.getlist('file_upload')
        
        
        # Check for existing order with same title
        order = ServiceOrder.objects.filter(user=request.user, title=title).first()
        
        is_new_order = False
        if not order:
            # Create new order if not exists
            order = ServiceOrder.objects.create(
                user=request.user,
                service_type=service_type,
                title=title,
                description=description,
                phone_number=phone_number,
                request_call=request_call,
            )
            is_new_order = True
        else:
             # Logic for Merged Order (Optional: Update description? For now, keep original to avoid overwrite)
             pass

        from .models import OrderFile
        files_added = 0
        
        for f in files:
            # Check for duplicate file name in this order
            if not OrderFile.objects.filter(order=order, original_filename=f.name).exists():
                 OrderFile.objects.create(order=order, file=f, file_type='source', original_filename=f.name)
                 files_added += 1

        if is_new_order:
            messages.success(request, "Order received successfully!")
        elif files_added > 0:
            messages.success(request, f"New files added to existing project: '{title}'")
        else:
             messages.info(request, f"Project '{title}' already exists and no new files were added.")
             
        return redirect('dashboard')
        
    return render(request, 'core/create_order.html')

@login_required
def profile_view(request):
    if request.method == 'POST':
        user = request.user
        if request.POST.get('full_name'):
            user.first_name = request.POST.get('full_name')
        if request.POST.get('email'):
            user.email = request.POST.get('email')
        user.save()

        profile, created = Profile.objects.get_or_create(user=user)
        profile.phone_number = request.POST.get('phone')
        
        if request.FILES.get('profile_image'):
            profile.profile_image = request.FILES.get('profile_image')
        
        profile.save()
        
        messages.success(request, "Profile updated successfully!")
        return redirect('profile')

    return render(request, 'core/profile.html')

@login_required
def payment_view(request, order_id):
    order = get_object_or_404(ServiceOrder, id=order_id, user=request.user)
    
    if request.method == 'POST':
        if request.FILES.get('payment_screenshot'):
            order.payment_screenshot = request.FILES.get('payment_screenshot')
            order.transaction_id = request.POST.get('transaction_id')
            order.is_paid = True
            order.save()
            
            messages.success(request, "Payment proof uploaded! We will verify shortly.")
            return redirect('dashboard')
            
    return render(request, 'core/payment.html', {'order': order})

# --- DYNAMIC PAGES LOGIC ---

def about_view(request):
    setting = SiteSetting.objects.first()
    team_members = TeamMember.objects.all().order_by('order')
    return render(request, 'core/about.html', {
        'setting': setting,
        'team_members': team_members
    })

def contact_view(request):
    setting = SiteSetting.objects.first()
    
    if request.method == 'POST':
        name = request.POST.get('name')
        email = request.POST.get('email')
        subject = request.POST.get('subject')
        message = request.POST.get('message')
        
        ContactMessage.objects.create(
            name=name, email=email, subject=subject, message=message
        )
        messages.success(request, "Message sent! We will contact you soon.")
        return redirect('contact')
        
    return render(request, 'core/contact.html', {'setting': setting})

def services_view(request):
    return render(request, 'core/services.html')

def faqs_view(request):
    return render(request, 'core/faqs.html')

def case_studies_view(request):
    case_studies = CaseStudy.objects.all().order_by('order')
    stats = AgencyStat.objects.all().order_by('order')
    setting = SiteSetting.objects.first()
    return render(request, 'core/case_studies.html', {
        'case_studies': case_studies,
        'agency_stats': stats,
        'setting': setting
    })

def terms_view(request):
    return render(request, 'core/terms.html')

def privacy(request):
    """Render the privacy policy page."""
    return render(request, 'core/privacy.html')

# --- NEW ORDER DETAIL & CHAT VIEW (UPDATED) ---
@login_required
def order_detail(request, order_id):
    order = get_object_or_404(ServiceOrder, id=order_id)
    
    # Security Check: Sirf apna order ya Admin dekh sake
    if request.user != order.user and not request.user.is_superuser and request.user.username != 'Hewor.order':
        messages.error(request, "You are not authorized to view this order.")
        return redirect('dashboard')
    
    if request.method == 'POST':
        # File Upload Logic
        if 'file_upload' in request.FILES:
            files = request.FILES.getlist('file_upload')
            upload_type = request.POST.get('upload_type', 'source')
            
            # Security: Non-admin users can ONLY upload source files
            if not request.user.is_superuser:
                upload_type = 'source'
            
            from .models import OrderFile
            for f in files:
                # Deduplication for additional uploads
                if not OrderFile.objects.filter(order=order, original_filename=f.name).exists():
                    OrderFile.objects.create(order=order, file=f, file_type=upload_type, original_filename=f.name)
            
            messages.success(request, f"{upload_type.title()} files uploaded successfully.")
            
            # Auto-update status if Admin delivers work
            if request.user.is_superuser and upload_type == 'delivery':
                if order.status != 'completed':
                    order.status = 'completed'
                    from django.utils import timezone
                    order.completed_at = timezone.now()
                    order.save()
            return redirect('order_detail', order_id=order.id)

        message = request.POST.get('message')
        
        if message:
            # Chat message create karo
            OrderChat.objects.create(
                order=order,
                sender=request.user,
                message=message
            )
            
            # Optional: Agar Admin reply kare toh status update ho
            if request.user.is_superuser:
                # Agar status pending hai to In Progress kar do
                if order.status == 'pending':
                    order.status = 'in_progress'
                    order.save()
            
            # Redirect to same page to avoid re-submission
            return redirect('order_detail', order_id=order.id)

    # Purani chats fetch karo (oldest first)
    chats = order.chats.all().order_by('created_at')
    
    # Separation of files
    source_files = order.files.filter(file_type='source')
    delivery_files = order.files.filter(file_type='delivery')
    
    return render(request, 'core/order_detail.html', {
        'order': order, 
        'chats': chats,
        'source_files': source_files,
        'delivery_files': delivery_files
    })

@login_required
def download_order_files(request, order_id, file_type):
    order = get_object_or_404(ServiceOrder, id=order_id)
    
    if request.user != order.user and not request.user.is_superuser:
        messages.error(request, "Access denied.")
        return redirect('dashboard')
    
    files = order.files.filter(file_type=file_type)
    
    if not files.exists():
        messages.error(request, "No files found.")
        return redirect('order_detail', order_id=order.id)
        
    zip_buffer = io.BytesIO()
    with zipfile.ZipFile(zip_buffer, 'w') as zip_file:
        for f in files:
            try:
                file_path = f.file.path
                file_name = f.file.name.split('/')[-1]
                zip_file.write(file_path, file_name)
            except FileNotFoundError:
                continue # Skip missing files
    
    zip_buffer.seek(0)
    response = HttpResponse(zip_buffer, content_type='application/zip')
    response['Content-Disposition'] = f'attachment; filename="{file_type}_files_{order.id}.zip"'
    return response

# --- CHATBOT API ---
from django.http import JsonResponse
from django.views.decorators.csrf import csrf_exempt
import json
from .chatbot_logic import get_chatbot_response

@csrf_exempt
def chatbot_api(request):
    if request.method == 'POST':
        try:
            data = json.loads(request.body)
            message = data.get('message', '')
            
            # Artificial Delay for realism (Optional, removed for speed)
            response_text = get_chatbot_response(request.user, message)
            
            return JsonResponse({'response': response_text, 'status': 'success'})
        except Exception as e:
            return JsonResponse({'response': "Error processing request.", 'status': 'error'})
    return JsonResponse({'status': 'invalid_method'})

# --- ORDER PANEL VIEWS ---
def order_panel_login(request):
    if request.user.is_authenticated:
        return redirect('order_panel_dashboard')
    
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            user = form.get_user()
            login(request, user)
            return redirect('order_panel_dashboard')
    else:
        form = AuthenticationForm()
    return render(request, 'core/order_panel_login.html', {'form': form})

@login_required(login_url='order_panel_login')
def order_panel_dashboard(request):
    orders = ServiceOrder.objects.all().order_by('-created_at')
    freelancers = Freelancer.objects.all().order_by('name')
    return render(request, 'core/order_panel_dashboard.html', {'orders': orders, 'freelancers': freelancers})

@login_required(login_url='order_panel_login')
def order_panel_upload(request, order_id):
    if request.method == 'POST':
        order = get_object_or_404(ServiceOrder, pk=order_id)
        files = request.FILES.getlist('file_upload')
        file_type = request.POST.get('file_type', 'delivery')
        
        for f in files:
            if not OrderFile.objects.filter(order=order, original_filename=f.name).exists():
                OrderFile.objects.create(
                    order=order, 
                    file=f, 
                    file_type=file_type, 
                    original_filename=f.name
                )
        messages.success(request, f"{len(files)} files uploaded.")
    return redirect('order_panel_dashboard')

@login_required(login_url='order_panel_login')
def order_panel_mark_delivered(request, order_id):
    if request.method == 'POST':
        order = get_object_or_404(ServiceOrder, pk=order_id)
        order.status = 'completed'
        from django.utils import timezone
        order.completed_at = timezone.now()
        order.save()
        messages.success(request, f"Order #{order.id} marked as Delivered.")
    return redirect('order_panel_dashboard')


@login_required(login_url='order_panel_login')
def order_panel_freelancers(request):
    if request.method == 'POST':
        name = request.POST.get('name')
        freelancer_id = request.POST.get('freelancer_id')
        phone = request.POST.get('phone')
        profession = request.POST.get('profession')
        address = request.POST.get('address')
        expertise = request.POST.get('expertise')
        profile_pic = request.FILES.get('profile_pic')
        password = request.POST.get('password') # Pass this from form
        
        qr_code = request.FILES.get('qr_code')
        payment_details = request.POST.get('payment_details')

        if Freelancer.objects.filter(freelancer_id=freelancer_id).exists():
            messages.error(request, f"Freelancer ID {freelancer_id} already exists.")
        else:
            try:
                # 1. Create User for Login
                user = User.objects.create_user(username=freelancer_id, password=password)
                
                # 2. Create Freelancer linked to User
                Freelancer.objects.create(
                    user=user,
                    name=name, freelancer_id=freelancer_id, phone=phone,
                    profession=profession, address=address, expertise=expertise,
                    profile_pic=profile_pic,
                    qr_code=qr_code, payment_details=payment_details
                )
                messages.success(request, f"Freelancer {name} added with Login ID: {freelancer_id}")
            except Exception as e:
                messages.error(request, f"Error creating freelancer: {str(e)}")
                
        return redirect('order_panel_freelancers')

    freelancers = Freelancer.objects.all().order_by('-joined_at')
    return render(request, 'core/order_panel_freelancers.html', {'freelancers': freelancers, 'active_tab': 'profiles'})

@login_required(login_url='order_panel_login')
def order_panel_freelancer_detail(request, freelancer_id):
    freelancer = get_object_or_404(Freelancer, id=freelancer_id)
    assigned_orders = ServiceOrder.objects.filter(freelancer=freelancer).order_by('-assigned_at')
    
    context = {
        'freelancer': freelancer,
        'assigned_orders': assigned_orders,
        'active_count': assigned_orders.filter(status='in_progress').count(),
        'completed_count': assigned_orders.filter(status='completed').count(),
        'pending_acceptance_count': assigned_orders.filter(freelancer_status='pending_acceptance').count(),
    }
    return render(request, 'core/admin/freelancer_details.html', context)

@login_required(login_url='order_panel_login')
def order_panel_assign_freelancer(request, order_id):
    if request.method == 'POST':
        order = get_object_or_404(ServiceOrder, pk=order_id)
        freelancer_id = request.POST.get('freelancer_id')
        roadmap = request.FILES.get('freelancer_roadmap')
        description = request.POST.get('freelancer_description')

        if freelancer_id:
            freelancer = get_object_or_404(Freelancer, id=freelancer_id)
            order.freelancer = freelancer
            
            if roadmap:
                order.freelancer_roadmap = roadmap
            if description:
                order.freelancer_description = description
           
            # Handle Payment Amount
            freelancer_payment = request.POST.get('freelancer_payment')
            if freelancer_payment:
                from decimal import Decimal, InvalidOperation
                try:
                    order.freelancer_payment = Decimal(freelancer_payment)
                except (ValueError, InvalidOperation):
                    messages.error(request, "Invalid payment amount.")
                    return redirect('order_panel_dashboard')
            
            # Set Assignment Details
            order.assigned_at = timezone.now()
            order.freelancer_status = 'pending_acceptance'
            # Default deadline (e.g., 2 days from now, can be made dynamic later)
            order.freelancer_deadline = timezone.now() + datetime.timedelta(days=2)
            
            order.save()
            
            # --- NOTIFICATION LOGIC ---
            # 1. Email Notification
            subject = f"New Work Assigned: {order.title}"
            message = f"""
            Hello {freelancer.name},
            
            You have been assigned a new project: "{order.title}".
            
            Price: Discussed with Admin.
            Description: {description or 'Check dashboard for details.'}
            
            Please login to your dashboard to Accept or Reject this order within 30 minutes.
            
            Login Here: {request.build_absolute_uri('/freelancer/login/')}
            """
            try:
                if freelancer.user.email:
                    send_mail(subject, message, settings.DEFAULT_FROM_EMAIL, [freelancer.user.email])
            except Exception as e:
                print(f"Email sending failed: {e}")
                
            # 2. WhatsApp Notification (Placeholder)
            print(f"--- WHATSAPP NOTIFICATION ---")
            print(f"To: {freelancer.phone}")
            print(f"Msg: New Order Assigned: {order.title}. Accept in 30 mins.")
            print(f"-----------------------------")
            
            messages.success(request, f"Order assigned to {freelancer.name}. Notifications sent.")
            
    return redirect('order_panel_dashboard')

@login_required(login_url='order_panel_login')
def order_panel_delete_freelancer(request, freelancer_id):
    if request.method == 'POST':
        freelancer = get_object_or_404(Freelancer, id=freelancer_id)
        freelancer.delete()
        messages.success(request, "Freelancer deleted successfully.")
    return redirect('order_panel_freelancers')

# --- 4b. FREELANCER ORDER CHAT ---
@login_required(login_url='order_panel_login')
def order_panel_freelancer_chat(request, order_id):
    # This view is for ADMIN side to chat with freelancer
    order = get_object_or_404(ServiceOrder, id=order_id)
    
    if request.method == 'POST':
        message = request.POST.get('message')
        attachment = request.FILES.get('attachment')
        if message or attachment:
            FreelancerChat.objects.create(
                order=order,
                sender=request.user,
                message=message,
                attachment=attachment
            )
            return redirect('order_panel_freelancer_chat', order_id=order.id)
            
    chats = order.freelancer_chats.all().order_by('created_at')
    return render(request, 'core/order_panel_chat.html', {'order': order, 'chats': chats})

# --- FREELANCER PORTAL VIEWS ---

def freelancer_login(request):
    if request.user.is_authenticated and hasattr(request.user, 'freelancer'):
        return redirect('freelancer_dashboard')
        
    if request.method == 'POST':
        form = AuthenticationForm(request, data=request.POST)
        if form.is_valid():
            user = form.get_user()
            # Ensure it's a freelancer account
            if hasattr(user, 'freelancer'):
                login(request, user)
                return redirect('freelancer_dashboard')
            else:
                messages.error(request, "Not a freelancer account.")
    else:
        form = AuthenticationForm()
    return render(request, 'core/freelancer_login.html', {'form': form})

@login_required(login_url='freelancer_login')
def freelancer_dashboard(request):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        logout(request)
        return redirect('freelancer_login')
        
    # --- AUTO-TIMEOUT LOGIC ---
    timeout_threshold = timezone.now() - datetime.timedelta(minutes=30)
    expired_orders = ServiceOrder.objects.filter(
        freelancer=freelancer,
        freelancer_status='pending_acceptance',
        assigned_at__lt=timeout_threshold
    )
    if expired_orders.exists():
        count = expired_orders.update(freelancer_status='timeout')
        print(f"Auto-timed out {count} orders for {freelancer.name}")

    # Get filter parameter
    filter_status = request.GET.get('filter', 'all')
    
    # Base query
    orders_query = ServiceOrder.objects.filter(freelancer=freelancer)
    
    # Apply filters
    if filter_status == 'pending':
        orders = orders_query.filter(freelancer_status='pending_acceptance')
    elif filter_status == 'active':
        orders = orders_query.filter(freelancer_status='accepted', status='in_progress')
    elif filter_status == 'due_soon':
        # Due within 3 days
        three_days_later = timezone.now() + datetime.timedelta(days=3)
        orders = orders_query.filter(
            freelancer_status='accepted',
            freelancer_deadline__lte=three_days_later,
            status='in_progress'
        )
    elif filter_status == 'completed':
        orders = orders_query.filter(status='completed')
    else:  # 'all'
        orders = orders_query
    
    orders = orders.order_by('-created_at')
    
    # --- CALCULATE STATS ---
    from django.db.models import Sum, Count, Q, Avg
    from datetime import timedelta
    
    # Total stats
    total_earned = orders_query.filter(is_freelancer_paid=True).count() * 1000  # Placeholder calculation
    active_count = orders_query.filter(freelancer_status='accepted', status='in_progress').count()
    completed_count = orders_query.filter(status='completed').count()
    pending_count = orders_query.filter(freelancer_status='pending_acceptance').count()
    
    # Earnings this week/month
    now = timezone.now()
    week_start = now - timedelta(days=now.weekday())
    month_start = now.replace(day=1, hour=0, minute=0, second=0, microsecond=0)
    
    week_earnings = orders_query.filter(
        is_freelancer_paid=True,
        created_at__gte=week_start
    ).count() * 1000  # Placeholder
    
    month_earnings = orders_query.filter(
        is_freelancer_paid=True,
        created_at__gte=month_start
    ).count() * 1000  # Placeholder
    
    # Due soon projects
    three_days_later = timezone.now() + datetime.timedelta(days=3)
    due_soon_count = orders_query.filter(
        freelancer_status='accepted',
        freelancer_deadline__lte=three_days_later,
        freelancer_deadline__gte=timezone.now(),
        status='in_progress'
    ).count()
    
    # --- NOTIFICATIONS ---
    notifications = []
    
    # New pending orders
    if pending_count > 0:
        notifications.append({
            'type': 'new_order',
            'icon': 'fa-inbox',
            'color': 'blue-500',
            'message': f'{pending_count} new order{"s" if pending_count > 1 else ""} waiting for acceptance',
            'time': 'Just now'
        })
    
    # Due soon alerts
    if due_soon_count > 0:
        notifications.append({
            'type': 'deadline',
            'icon': 'fa-clock',
            'color': 'yellow-500',
            'message': f'{due_soon_count} project{"s" if due_soon_count > 1 else ""} due within 3 days',
            'time': 'Today'
        })
    
    # Payment received (recent)
    recent_paid = orders_query.filter(
        is_freelancer_paid=True,
        created_at__gte=now - timedelta(days=7)
    ).count()
    if recent_paid > 0:
        notifications.append({
            'type': 'payment',
            'icon': 'fa-money-bill-wave',
            'color': 'green-500',
            'message': f'Payment received for {recent_paid} project{"s" if recent_paid > 1 else ""}',
            'time': 'This week'
        })
    
    # --- QUICK ACTIONS ---
    quick_actions = []
    
    # Pending acceptance
    if pending_count > 0:
        quick_actions.append({
            'icon': 'fa-check-circle',
            'color': 'green',
            'text': f'Accept {pending_count} new order{"s" if pending_count > 1 else ""}',
            'link': '?filter=pending'
        })
    
    # Active projects
    if active_count > 0:
        quick_actions.append({
            'icon': 'fa-tasks',
            'color': 'blue',
            'text': f'Continue {active_count} active project{"s" if active_count > 1 else ""}',
            'link': '?filter=active'
        })
    
    # Due soon
    if due_soon_count > 0:
        quick_actions.append({
            'icon': 'fa-exclamation-triangle',
            'color': 'yellow',
            'text': f'Check {due_soon_count} urgent deadline{"s" if due_soon_count > 1 else ""}',
            'link': '?filter=due_soon'
        })
    
    context = {
        'orders': orders,
        'filter_status': filter_status,
        'stats': {
            'total_earned': total_earned,
            'week_earnings': week_earnings,
            'month_earnings': month_earnings,
            'active_count': active_count,
            'completed_count': completed_count,
            'pending_count': pending_count,
            'due_soon_count': due_soon_count,
            'rating': 4.8,  # Placeholder - implement rating system later
        },
        'notifications': notifications[:5],  # Show max 5
        'quick_actions': quick_actions,
    }
    
    return render(request, 'core/freelancer_dashboard.html', context)

@login_required(login_url='order_panel_login')
def order_panel_pay_freelancer(request, order_id):
    order = get_object_or_404(ServiceOrder, pk=order_id)
    if request.method == 'POST':
        try:
            transaction_id = request.POST.get('transaction_id')
            screenshot = request.FILES.get('payment_screenshot')
            
            if transaction_id:
                order.is_freelancer_paid = True
                order.freelancer_transaction_id = transaction_id
                if screenshot:
                    order.freelancer_payment_screenshot = screenshot
                order.save()
                messages.success(request, f"Payment recorded for Freelancer: {order.freelancer.name}")
            else:
                messages.error(request, "Transaction ID is required.")
        except Exception as e:
            messages.error(request, f"Error processing payment: {str(e)}")
            
    return redirect('manage_freelancer_works')

# Redoing the tool call with the actual code after I update the model. 
# actually, I should update the model FIRST.
# Cancelling this replacement to update model first.


@login_required(login_url='order_panel_login')
def manage_freelancer_works(request):
    # Fetch all orders that have been assigned to a freelancer
    assignments = ServiceOrder.objects.filter(freelancer__isnull=False).order_by('-assigned_at')
    
    context = {
        'assignments': assignments,
        'active_tab': 'works'
    }
    return render(request, 'core/manage_freelancer_works.html', context)


@login_required(login_url='freelancer_login')
def freelancer_order_detail(request, order_id):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        logout(request)
        return redirect('freelancer_login')
        
    order = get_object_or_404(ServiceOrder, id=order_id, freelancer=freelancer)
    
    if request.method == 'POST':
        action = request.POST.get('action')
        
        if action == 'upload_work':
            files = request.FILES.getlist('files')
            if files:
                for f in files:
                    OrderFile.objects.create(
                        order=order,
                        file=f,
                        file_type='freelancer_upload',
                        original_filename=f.name
                    )
                messages.success(request, f"{len(files)} files uploaded successfully.")
                
        elif action == 'send_message':
            message = request.POST.get('message')
            attachment = request.FILES.get('attachment')
            if message or attachment:
                FreelancerChat.objects.create(
                    order=order,
                    sender=request.user,
                    message=message,
                    attachment=attachment
                )
        return redirect('freelancer_order_detail', order_id=order.id)
            
    chats = order.freelancer_chats.all().order_by('created_at')
    uploaded_files = order.files.filter(file_type='freelancer_upload').order_by('-uploaded_at')
    client_files = order.files.filter(file_type='source').order_by('-uploaded_at')  # Client's original files
    
    return render(request, 'core/freelancer_order_detail.html', {
        'order': order,
        'chats': chats,
        'uploaded_files': uploaded_files,
        'client_files': client_files,  # Pass client files to template
    })

@login_required(login_url='freelancer_login')
def freelancer_accept_order(request, order_id):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        return redirect('freelancer_login')
        
    order = get_object_or_404(ServiceOrder, id=order_id, freelancer=freelancer)
    
    if order.freelancer_status != 'pending_acceptance':
        messages.error(request, "Order is not pending acceptance.")
        return redirect('freelancer_dashboard')
    
    # Accept the order
    order.freelancer_status = 'accepted'
    order.status = 'in_progress'  # Update main order status
    order.save()
    
    messages.success(request, f"You've successfully accepted the order: {order.title}")
    return redirect('freelancer_order_detail', order_id=order.id)

def freelancer_reject_order(request, order_id):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        return redirect('freelancer_login')
        
    order = get_object_or_404(ServiceOrder, id=order_id, freelancer=freelancer)
    
    if order.freelancer_status != 'pending_acceptance':
        messages.error(request, "Order is not pending acceptance.")
        return redirect('freelancer_dashboard')
    
    # Reject the order and unassign freelancer
    order.freelancer_status = 'unassigned'
    order.freelancer = None
    order.save()
    
    messages.info(request, f"You've rejected the order: {order.title}")
    return redirect('freelancer_dashboard')

@login_required(login_url='freelancer_login')
def freelancer_profile(request):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        logout(request)
        return redirect('freelancer_login')
    
    if request.method == 'POST':
        # Update profile
        freelancer.name = request.POST.get('name', freelancer.name)
        freelancer.phone = request.POST.get('phone', freelancer.phone)
        freelancer.email = request.POST.get('email', freelancer.email)
        freelancer.save()
        
        messages.success(request, "Profile updated successfully!")
        return redirect('freelancer_profile')
    
    return render(request, 'core/freelancer_profile.html', {'freelancer': freelancer})

# --- FREE TOOLS ---

def merge_pdf_tool(request):
    """
    View to handle Free PDF Merge tool.
    Limits: Max 200MB total payload.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        
        if not files:
            messages.error(request, "Please select at least one PDF file.")
            return redirect('merge_pdf_tool')

        # 1. Size Validation (Max 200MB)
        MAX_SIZE_MB = 100
        total_size = sum(f.size for f in files)
        if total_size > MAX_SIZE_MB * 1024 * 1024:
            messages.error(request, f"Total file size exceeds {MAX_SIZE_MB}MB limit. Your files: {round(total_size / (1024*1024), 2)}MB")
            return redirect('merge_pdf_tool')

        try:
            # 2. Merge Logic using PyMuPDF (fitz)
            merged_pdf = fitz.open()

            for f in files:
                # Read file from memory
                file_stream = f.read()
                # Open PDF from memory stream
                with fitz.open(stream=file_stream, filetype="pdf") as pdf_doc:
                    merged_pdf.insert_pdf(pdf_doc)

            # 3. Return Response
            pdf_bytes = merged_pdf.tobytes()
            merged_pdf.close()

            response = HttpResponse(pdf_bytes, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="hewor_merged.pdf"'
            return response

        except Exception as e:
            logger.error(f"Error merging PDFs: {e}")
            messages.error(request, f"Error processing files: {str(e)}")
            return redirect('merge_pdf_tool')

    return render(request, 'core/merge_pdf.html')

def split_pdf_tool(request):
    """
    View to handle Free Split PDF tool.
    Supports single or multiple files.
    Input: PDF file(s) + comma-separated page numbers.
    Output: ZIP file of split parts.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        split_pages_str = request.POST.get('split_pages', '')

        if not files:
            messages.error(request, "Please upload at least one PDF file.")
            return redirect('split_pdf_tool')

        try:
            # Parse split pages
            split_at_pages = [int(p.strip()) for p in split_pages_str.split(',') if p.strip().isdigit()]
            split_at_pages.sort()
            
            if not split_at_pages:
                messages.error(request, "Please enter valid page numbers.")
                return redirect('split_pdf_tool')

            # Prepare ZIP buffer for output
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                
                for file_idx, file in enumerate(files):
                    # Check Limit
                    MAX_SIZE_MB = 100
                    if file.size > MAX_SIZE_MB * 1024 * 1024:
                         messages.error(request, f"File {file.name} exceeds limit.")
                         return redirect('split_pdf_tool')

                    doc = fitz.open(stream=file.read(), filetype="pdf")
                    total_pages = doc.page_count
                    
                    # Determine split ranges
                    start_page = 0
                    part_num = 1
                    
                    # Logic to handle splits
                    # e.g. split at 5, 10 -> [0-4], [5-9], [10-end]
                    
                    current_splits = [p for p in split_at_pages if p < total_pages]
                    
                    # Logic: If split page is 5, it means split AFTER page 5 (pages 1-5 in output).
                    # fitz uses 0-indexed. So split point 5 means page index 5 is start of NEXT doc.
                    # Range 1: 0 to 4.
                    
                    ranges = []
                    last_split = 0
                    
                    for split_point in current_splits:
                        # User enters "5" -> they probably mean page 5 (visual). 
                        # Code usually interprets "split after page 5".
                        # So range is 0 to 4 (5 pages).
                        # p index = 5 is the 6th page.
                        # Wait, let's stick to standard PDF logic:
                        # Input: "5" -> File 1: Pages 1-5. File 2: Pages 6-end.
                        # Index range: [0, 5) and [5, total)
                        
                        ranges.append((last_split, split_point))
                        last_split = split_point
                        
                    # Add final range
                    if last_split < total_pages:
                        ranges.append((last_split, total_pages))
                        
                    base_name = os.path.splitext(file.name)[0]
                    
                    for r_start, r_end in ranges:
                        if r_start >= r_end: continue
                        
                        new_doc = fitz.open()
                        new_doc.insert_pdf(doc, from_page=r_start, to_page=r_end-1)
                        
                        # Save to ZIP
                        out_pdf_bytes = new_doc.write()
                        zip_file.writestr(f"{base_name}_part_{part_num}.pdf", out_pdf_bytes)
                        
                        new_doc.close()
                        part_num += 1
                        
                    doc.close()

            zip_buffer.seek(0)
            response = HttpResponse(zip_buffer, content_type='application/zip')
            response['Content-Disposition'] = 'attachment; filename="hewor_split_files.zip"'
            return response

        except Exception as e:
            logger.error(f"Error splitting PDF: {e}")
            messages.error(request, f"Error processing file: {str(e)}")
            return redirect('split_pdf_tool')

    return render(request, 'core/split_pdf.html')

def compress_pdf_tool(request):
    """
    View to handle Free Compress PDF tool.
    Supports batch processing.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        
        if not files:
            messages.error(request, "Please upload at least one PDF file.")
            return redirect('compress_pdf_tool')

        MAX_SIZE_MB = 100
        # Check total size logic if desired, or per file. 
        # Using per file for now or simple sum.
        
        try:
            # If single file -> return PDF
            if len(files) == 1:
                file = files[0]
                if file.size > MAX_SIZE_MB * 1024 * 1024:
                    messages.error(request, f"File size exceeds {MAX_SIZE_MB}MB limit.")
                    return redirect('compress_pdf_tool')
                
                doc = fitz.open(stream=file.read(), filetype="pdf")
                # Compress
                # garbage=4 (deduplicate), deflate=True (compress streams)
                out_bytes = doc.write(garbage=4, deflate=True)
                doc.close()
                
                response = HttpResponse(out_bytes, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="compressed_{file.name}"'
                return response
            
            # If multiple files -> return ZIP
            else:
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    for file in files:
                        doc = fitz.open(stream=file.read(), filetype="pdf")
                        out_bytes = doc.write(garbage=4, deflate=True)
                        zip_file.writestr(f"compressed_{file.name}", out_bytes)
                        doc.close()
                
                zip_buffer.seek(0)
                response = HttpResponse(zip_buffer, content_type='application/zip')
                response['Content-Disposition'] = 'attachment; filename="hewor_compressed_batch.zip"'
                return response

        except Exception as e:
            logger.error(f"Error compressing PDF: {e}")
            messages.error(request, f"Error processing file: {str(e)}")
            return redirect('compress_pdf_tool')

    return render(request, 'core/compress_pdf.html')

def pdf_to_word_tool(request):
    """
    View to handle Free PDF to Word tool.
    Note: pdf2docx library removed due to deployment issues.
    This tool is temporarily disabled.
    """
    if request.method == 'POST':
        messages.error(request, "PDF to Word conversion is temporarily unavailable. Please try Word to PDF instead!")
        return redirect('pdf_to_word_tool')

    return render(request, 'core/pdf_to_word.html')
        



def pdf_to_ppt_tool(request):
    """
    View to handle Free PDF to PowerPoint tool.
    Uses PyMuPDF (fitz) to render pages as images and python-pptx to create slides.
    Supports batch processing.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('pdf_to_ppt_tool')

        MAX_SIZE_MB = 100
        temp_files_to_clean = []
        
        try:
            def convert_single_pdf_to_pptx(pdf_file, output_pptx_path):
                # Save input PDF to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                    for chunk in pdf_file.chunks():
                        tmp_pdf.write(chunk)
                    temp_pdf_path = tmp_pdf.name
                    temp_files_to_clean.append(temp_pdf_path)

                doc = fitz.open(temp_pdf_path)
                prs = Presentation()
                
                # Standard slide size (10x7.5 inches) or Wide (13.33x7.5)
                # We can try to match PDF aspect ratio or just use standard
                # Let's set slide size to match the first page of PDF if possible
                if len(doc) > 0:
                    page = doc[0]
                    # PyMuPDF uses points. 1 inch = 72 points.
                    # python-pptx uses EMU. Inches(1)
                    width_pt = page.rect.width
                    height_pt = page.rect.height
                    
                    prs.slide_width = int(width_pt * 914400 / 72)
                    prs.slide_height = int(height_pt * 914400 / 72)

                blank_slide_layout = prs.slide_layouts[6] 

                for page_num in range(len(doc)):
                    page = doc[page_num]
                    # Render high quality image (zoom=2)
                    mat = fitz.Matrix(2, 2)
                    pix = page.get_pixmap(matrix=mat)
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_img:
                        pix.save(tmp_img.name)
                        tmp_img_path = tmp_img.name
                        temp_files_to_clean.append(tmp_img_path)
                    
                    slide = prs.slides.add_slide(blank_slide_layout)
                    slide.shapes.add_picture(tmp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

                prs.save(output_pptx_path)
                doc.close()

            # Single File Case
            if len(files) == 1:
                file = files[0]
                if file.size > MAX_SIZE_MB * 1024 * 1024:
                     messages.error(request, f"File size exceeds {MAX_SIZE_MB}MB limit.")
                     return redirect('pdf_to_ppt_tool')

                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx:
                    temp_pptx_path = tmp_pptx.name
                    temp_files_to_clean.append(temp_pptx_path)
                
                convert_single_pdf_to_pptx(file, temp_pptx_path)

                with open(temp_pptx_path, 'rb') as pptx_file:
                    response = HttpResponse(pptx_file.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    base_name = os.path.splitext(file.name)[0]
                    response['Content-Disposition'] = f'attachment; filename="{base_name}.pptx"'
                
                return response
            
            # Batch Case
            else:
                zip_buffer = io.BytesIO()
                with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zip_file:
                    for file in files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_pptx:
                            temp_pptx_path = tmp_pptx.name
                            temp_files_to_clean.append(temp_pptx_path)
                        
                        try:
                            convert_single_pdf_to_pptx(file, temp_pptx_path)
                            base_name = os.path.splitext(file.name)[0]
                            zip_file.write(temp_pptx_path, f"{base_name}.pptx")
                        except Exception as sub_e:
                            logger.error(f"Error in batch pdf2pptx for {file.name}: {sub_e}")
                            pass

                zip_buffer.seek(0)
                response = HttpResponse(zip_buffer, content_type='application/zip')
                response['Content-Disposition'] = 'attachment; filename="hewor_converted_ppt_batch.zip"'
                return response

        except Exception as e:
            logger.error(f"Error converting PDF to PPTX: {e}")
            messages.error(request, f"Error processing file: {str(e)}")
            return redirect('pdf_to_ppt_tool')
        
        finally:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    try:
                        os.remove(path)
                    except:
                        pass

    return render(request, 'core/pdf_to_powerpoint.html')


@login_required(login_url='freelancer_login')
def freelancer_reject_order(request, order_id):
    try:
        freelancer = request.user.freelancer
    except Freelancer.DoesNotExist:
        return redirect('freelancer_login')
        
    order = get_object_or_404(ServiceOrder, id=order_id, freelancer=freelancer)
    
    if order.freelancer_status == 'pending_acceptance':
        order.freelancer_status = 'rejected'
        order.save()
        messages.info(request, "Order rejected.")
        
    return redirect('freelancer_dashboard')

def pdf_to_excel_tool(request):
    """
    View to handle Free PDF to Excel tool.
    Uses pdfplumber to extract tables and pandas to save as Excel.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('pdf_to_excel_tool')

        temp_files_to_clean = []
        
        try:
            def convert_single_pdf_to_excel(pdf_file, output_excel_path):
                # Save input PDF to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                    for chunk in pdf_file.chunks():
                        tmp_pdf.write(chunk)
                    temp_pdf_path = tmp_pdf.name
                    temp_files_to_clean.append(temp_pdf_path)

                # Extract and write to Excel
                has_tables = False
                with pdfplumber.open(temp_pdf_path) as pdf:
                    with pd.ExcelWriter(output_excel_path, engine='openpyxl') as writer:
                        for i, page in enumerate(pdf.pages):
                            tables = page.extract_tables()
                            for j, table in enumerate(tables):
                                if table:
                                    has_tables = True
                                    # Create DataFrame. defaulting to no header to preserve exact structure
                                    df = pd.DataFrame(table)
                                    sheet_name = f'Page_{i+1}_Table_{j+1}'
                                    # Sheet name limit is 31 chars
                                    if len(sheet_name) > 31:
                                        sheet_name = sheet_name[:31]
                                    df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)
                
                if not has_tables:
                    # Create a dummy sheet if no tables found to avoid error
                    with pd.ExcelWriter(output_excel_path, engine='openpyxl') as writer:
                        pd.DataFrame(["No tables found in this PDF"]).to_excel(writer, sheet_name="Info", index=False, header=False)

            results = []
            
            if len(files) == 1:
                # Single file case
                uploaded_file = files[0]
                with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_out:
                    output_path = tmp_out.name
                    temp_files_to_clean.append(output_path)
                
                convert_single_pdf_to_excel(uploaded_file, output_path)
                
                # Serve file
                with open(output_path, 'rb') as f:
                    file_data = f.read()
                    
                response = HttpResponse(file_data, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".pdf", "")}.xlsx"'
                
                # Cleanup
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                        
                return response
            
            else:
                # Batch processing
                zip_filename = "hewor_converted_excel_files.zip"
                with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp_zip:
                    zip_path = tmp_zip.name
                    temp_files_to_clean.append(zip_path)
                
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for uploaded_file in files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_out:
                            output_path = tmp_out.name
                            temp_files_to_clean.append(output_path)
                        
                        convert_single_pdf_to_excel(uploaded_file, output_path)
                        zipf.write(output_path, arcname=f"{uploaded_file.name.replace('.pdf', '')}.xlsx")
                
                with open(zip_path, 'rb') as f:
                    zip_data = f.read()
                    
                response = HttpResponse(zip_data, content_type='application/zip')
                response['Content-Disposition'] = f'attachment; filename="{zip_filename}"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                        
                return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('pdf_to_excel_tool')

    return render(request, 'core/pdf_to_excel.html')

def word_to_pdf_tool(request):
    """
    View to handle Free Word to PDF tool.
    Tries docx2pdf first (best quality), falls back to python-docx + xhtml2pdf.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('word_files')
        
        if not files:
            messages.error(request, "Please upload a Word file.")
            return redirect('word_to_pdf_tool')

        temp_files_to_clean = []
        
        try:
            def convert_single_word_to_pdf(docx_file, output_pdf_path):
                # Save input DOCX to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_docx:
                    for chunk in docx_file.chunks():
                        tmp_docx.write(chunk)
                    temp_docx_path = tmp_docx.name
                    temp_files_to_clean.append(temp_docx_path)

                # Conversion Logic
                conversion_success = False
                
                # Method 1: docx2pdf (Requires MS Word installed)
                # We skip this for now on server environments usually, but let's try safely
                try:
                    from docx2pdf import convert
                    # Note: docx2pdf might fail if no Word is installed or on Linux headless
                    # convert(temp_docx_path, output_pdf_path) 
                    # Commented out to prefer the safe method for now to avoid hanging
                    pass
                except:
                    pass

                if not conversion_success:
                    # Method 2: python-docx -> HTML -> xhtml2pdf
                    try:
                        import docx
                        from xhtml2pdf import pisa
                        
                        doc = docx.Document(temp_docx_path)
                        html_content = "<html><head><style>body { font-family: Helvetica, sans-serif; }</style></head><body>"
                        
                        for para in doc.paragraphs:
                            # Basic styling handling
                            style = para.style.name.lower()
                            text = para.text.strip()
                            if not text:
                                continue
                                
                            if 'heading' in style:
                                level = style.replace('heading ', '')
                                html_content += f"<h{level}>{text}</h{level}>"
                            else:
                                html_content += f"<p>{text}</p>"
                                
                        html_content += "</body></html>"
                        
                        with open(output_pdf_path, "wb") as pdf_file:
                            pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)
                            
                        if pisa_status.err:
                            raise Exception("PDF generation failed")
                        
                        conversion_success = True
                    except Exception as e:
                        logger.error(f"Word conversion failed: {e}")
                        raise e

            # Batch processing logic
            if len(files) == 1:
                uploaded_file = files[0]
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                    output_path = tmp_out.name
                    temp_files_to_clean.append(output_path)
                
                convert_single_word_to_pdf(uploaded_file, output_path)
                
                with open(output_path, 'rb') as f:
                    file_data = f.read()
                    
                response = HttpResponse(file_data, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".docx", "").replace(".doc", "")}.pdf"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response
            
            else:
                zip_filename = "hewor_converted_pdf_files.zip"
                with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp_zip:
                    zip_path = tmp_zip.name
                    temp_files_to_clean.append(zip_path)
                
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for uploaded_file in files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                            output_path = tmp_out.name
                            temp_files_to_clean.append(output_path)
                        
                        try:
                            convert_single_word_to_pdf(uploaded_file, output_path)
                            zipf.write(output_path, arcname=f"{uploaded_file.name.replace('.docx', '').replace('.doc', '')}.pdf")
                        except Exception as e:
                             logger.error(f"Failed to convert {uploaded_file.name}: {e}")
                
                with open(zip_path, 'rb') as f:
                    zip_data = f.read()
                    
                response = HttpResponse(zip_data, content_type='application/zip')
                response['Content-Disposition'] = f'attachment; filename="{zip_filename}"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('word_to_pdf_tool')

    return render(request, 'core/word_to_pdf.html')

def excel_to_pdf_tool(request):
    """
    View to handle Free Excel to PDF tool.
    Uses pandas to read Excel and xhtml2pdf to generate PDF.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('excel_files')
        
        if not files:
            messages.error(request, "Please upload an Excel file.")
            return redirect('excel_to_pdf_tool')

        temp_files_to_clean = []
        
        try:
            def convert_single_excel_to_pdf(xlsx_file, output_pdf_path):
                # Save input XLSX to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as tmp_xlsx:
                    for chunk in xlsx_file.chunks():
                        tmp_xlsx.write(chunk)
                    temp_xlsx_path = tmp_xlsx.name
                    temp_files_to_clean.append(temp_xlsx_path)

                try:
                    import pandas as pd
                    from xhtml2pdf import pisa
                    
                    # Read all sheets
                    xls = pd.ExcelFile(temp_xlsx_path)
                    html_content = "<html><head><style>"
                    html_content += "@page { size: A4 landscape; margin: 1cm; }"
                    html_content += "table { border-collapse: collapse; width: 100%; margin-bottom: 20px; font-size: 10px; }"
                    html_content += "th, td { border: 1px solid #ddd; padding: 5px; text-align: left; }"
                    html_content += "th { background-color: #f2f2f2; font-weight: bold; }"
                    html_content += "h2 { font-family: sans-serif; color: #333; }"
                    html_content += "</style></head><body>"
                    
                    for sheet_name in xls.sheet_names:
                        df = pd.read_excel(xls, sheet_name=sheet_name)
                        if not df.empty:
                            html_content += f"<h2>Sheet: {sheet_name}</h2>"
                            # Convert to HTML
                            table_html = df.to_html(index=False, classes='table')
                            html_content += table_html
                            html_content += "<br/>"
                            
                    html_content += "</body></html>"
                    
                    with open(output_pdf_path, "wb") as pdf_file:
                        pisa_status = pisa.CreatePDF(html_content, dest=pdf_file)
                        
                    if pisa_status.err:
                        raise Exception("PDF generation failed")

                except Exception as e:
                    logger.error(f"Excel conversion failed: {e}")
                    raise e

            # Batch processing logic
            if len(files) == 1:
                uploaded_file = files[0]
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                    output_path = tmp_out.name
                    temp_files_to_clean.append(output_path)
                
                convert_single_excel_to_pdf(uploaded_file, output_path)
                
                with open(output_path, 'rb') as f:
                    file_data = f.read()
                    
                response = HttpResponse(file_data, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".xlsx", "").replace(".xls", "")}.pdf"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response
            
            else:
                zip_filename = "hewor_converted_pdf_files.zip"
                with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp_zip:
                    zip_path = tmp_zip.name
                    temp_files_to_clean.append(zip_path)
                
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for uploaded_file in files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                            output_path = tmp_out.name
                            temp_files_to_clean.append(output_path)
                        
                        try:
                            convert_single_excel_to_pdf(uploaded_file, output_path)
                            zipf.write(output_path, arcname=f"{uploaded_file.name.replace('.xlsx', '').replace('.xls', '')}.pdf")
                        except Exception as e:
                             logger.error(f"Failed to convert {uploaded_file.name}: {e}")
                
                with open(zip_path, 'rb') as f:
                    zip_data = f.read()
                    
                response = HttpResponse(zip_data, content_type='application/zip')
                response['Content-Disposition'] = f'attachment; filename="{zip_filename}"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('excel_to_pdf_tool')

    return render(request, 'core/excel_to_pdf.html')

def ppt_to_pdf_tool(request):
    """
    View to handle Free PowerPoint to PDF tool.
    Extracts text content from slides and generates a PDF report.
    Note: Does not preserve layout (requires LibreOffice for that).
    """
    if request.method == 'POST':
        files = request.FILES.getlist('ppt_files')
        
        if not files:
            messages.error(request, "Please upload a PowerPoint file.")
            return redirect('ppt_to_pdf_tool')

        temp_files_to_clean = []
        
        try:
            def convert_single_ppt_to_pdf(ppt_file, output_pdf_path):
                # Save input PPTX to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_ppt:
                    for chunk in ppt_file.chunks():
                        tmp_ppt.write(chunk)
                    temp_ppt_path = tmp_ppt.name
                    temp_files_to_clean.append(temp_ppt_path)

                try:
                    from pptx import Presentation
                    from reportlab.lib.pagesizes import letter
                    from reportlab.pdfgen import canvas
                    from reportlab.lib import colors
                    from reportlab.lib.styles import getSampleStyleSheet
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
                    
                    # Use Platypus for better text wrapping
                    doc = SimpleDocTemplate(output_pdf_path, pagesize=letter)
                    story = []
                    styles = getSampleStyleSheet()
                    
                    prs = Presentation(temp_ppt_path)
                    
                    for i, slide in enumerate(prs.slides):
                        # Slide Header
                        story.append(Paragraph(f"<b>Slide {i+1}</b>", styles['Heading2']))
                        story.append(Spacer(1, 12))
                        
                        # Extract Text
                        slide_text = []
                        # shapes are not always in reading order, but we try
                        for shape in slide.shapes:
                            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                                for paragraph in shape.text_frame.paragraphs:
                                    text = paragraph.text.strip()
                                    if text:
                                        story.append(Paragraph(text, styles['BodyText']))
                                        story.append(Spacer(1, 6))
                        
                        story.append(PageBreak())
                        
                    doc.build(story)

                except Exception as e:
                    logger.error(f"PPT conversion failed: {e}")
                    raise e

            # Batch processing logic
            if len(files) == 1:
                uploaded_file = files[0]
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                    output_path = tmp_out.name
                    temp_files_to_clean.append(output_path)
                
                convert_single_ppt_to_pdf(uploaded_file, output_path)
                
                with open(output_path, 'rb') as f:
                    file_data = f.read()
                    
                response = HttpResponse(file_data, content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename="{uploaded_file.name.replace(".pptx", "").replace(".ppt", "")}.pdf"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response
            
            else:
                zip_filename = "hewor_converted_pdf_files.zip"
                with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp_zip:
                    zip_path = tmp_zip.name
                    temp_files_to_clean.append(zip_path)
                
                with zipfile.ZipFile(zip_path, 'w') as zipf:
                    for uploaded_file in files:
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                            output_path = tmp_out.name
                            temp_files_to_clean.append(output_path)
                        
                        try:
                            convert_single_ppt_to_pdf(uploaded_file, output_path)
                            zipf.write(output_path, arcname=f"{uploaded_file.name.replace('.pptx', '').replace('.ppt', '')}.pdf")
                        except Exception as e:
                             logger.error(f"Failed to convert {uploaded_file.name}: {e}")
                
                with open(zip_path, 'rb') as f:
                    zip_data = f.read()
                    
                response = HttpResponse(zip_data, content_type='application/zip')
                response['Content-Disposition'] = f'attachment; filename="{zip_filename}"'
                
                for path in temp_files_to_clean:
                    if os.path.exists(path):
                        os.remove(path)
                return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('ppt_to_pdf_tool')

    return render(request, 'core/ppt_to_pdf.html')

def pdf_to_jpg_tool(request):
    """
    View to handle Free PDF to JPG tool.
    Converts PDF pages to high-quality JPG images.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('pdf_to_jpg_tool')

        temp_files_to_clean = []
        
        try:
            def convert_single_pdf_to_jpgs(pdf_file, output_zip_path):
                # Save input PDF to temp file
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                    for chunk in pdf_file.chunks():
                        tmp_pdf.write(chunk)
                    temp_pdf_path = tmp_pdf.name
                    temp_files_to_clean.append(temp_pdf_path)

                # Open PDF and convert pages
                doc = fitz.open(temp_pdf_path)
                base_name = pdf_file.name.replace('.pdf', '')
                
                with zipfile.ZipFile(output_zip_path, 'w') as zipf:
                    for i, page in enumerate(doc):
                        # clear resolution (zoom=2 means 144dpi approx, good for screen)
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                        
                        # Save each page as JPG
                        img_filename = f"{base_name}_page_{i+1}.jpg"
                        with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_img:
                            pix.save(tmp_img.name)
                            temp_files_to_clean.append(tmp_img.name)
                            
                            # Add to zip
                            zipf.write(tmp_img.name, arcname=img_filename)
            
            # Use a main ZIP for the download
            final_zip_filename = "hewor_converted_jpgs.zip"
            with tempfile.NamedTemporaryFile(delete=False, suffix='.zip') as tmp_final_zip:
                final_zip_path = tmp_final_zip.name
                temp_files_to_clean.append(final_zip_path)
            
            # Since even a single PDF produces multiple JPGs, we always return a ZIP of JPGs (or ZIP of folders if batch)
            # Simpler approach: Just put all images in one ZIP. If batch, prefix with filename.
            
            with zipfile.ZipFile(final_zip_path, 'w') as final_zip:
                for uploaded_file in files:
                    # Save input PDF to temp
                    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                        for chunk in uploaded_file.chunks():
                            tmp_pdf.write(chunk)
                        temp_pdf_path = tmp_pdf.name
                        temp_files_to_clean.append(temp_pdf_path)
                    
                    doc = fitz.open(temp_pdf_path)
                    base_name = uploaded_file.name.replace('.pdf', '')
                    
                    for i, page in enumerate(doc):
                         pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                         img_filename = f"{base_name}_page_{i+1}.jpg"
                         
                         with tempfile.NamedTemporaryFile(delete=False, suffix='.jpg') as tmp_img:
                            pix.save(tmp_img.name)
                            temp_files_to_clean.append(tmp_img.name)
                            final_zip.write(tmp_img.name, arcname=img_filename)

            with open(final_zip_path, 'rb') as f:
                zip_data = f.read()
                
            response = HttpResponse(zip_data, content_type='application/zip')
            response['Content-Disposition'] = f'attachment; filename="{final_zip_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('pdf_to_jpg_tool')

    return render(request, 'core/pdf_to_jpg.html')

def jpg_to_pdf_tool(request):
    """
    View to handle Free JPG to PDF tool.
    Converts uploaded images to a single PDF.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('jpg_files')
        
        if not files:
            messages.error(request, "Please upload image files.")
            return redirect('jpg_to_pdf_tool')

        temp_files_to_clean = []
        
        try:
            # Save all images first
            image_paths = []
            for img_file in files:
                suffix = os.path.splitext(img_file.name)[1].lower()
                if suffix not in ['.jpg', '.jpeg', '.png']:
                    # Force valid suffix if missing or allow conversion later?
                    # img2pdf requires jpg/jpeg naming or valid header
                    suffix = '.jpg'
                    
                with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_img:
                    for chunk in img_file.chunks():
                        tmp_img.write(chunk)
                    tmp_img.flush()
                    image_paths.append(tmp_img.name)
                    temp_files_to_clean.append(tmp_img.name)

            # Define Output PDF
            final_filename = "hewor_images_combined.pdf"
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                output_pdf_path = tmp_pdf.name
                temp_files_to_clean.append(output_pdf_path)

            # TRY img2pdf first (best quality, direct embedding)
            try:
                import img2pdf
                # img2pdf might fail on some PNGs (alpha channel), so fallback to PIL is safer for generic "Image to PDF".
                # But let's try it if inputs are strictly JPEG.
                # Actually, PIL is robust for user uploads that might be messy.
                # Let's use PIL for maximum compatibility.
                
                from PIL import Image
                
                pil_images = []
                for p in image_paths:
                    img = Image.open(p)
                    # Convert to RGB to avoid mode errors (e.g. RGBA -> PDF issue)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    pil_images.append(img)
                
                if pil_images:
                    base_image = pil_images[0]
                    other_images = pil_images[1:] if len(pil_images) > 1 else []
                    
                    base_image.save(output_pdf_path, "PDF", resolution=100.0, save_all=True, append_images=other_images)
                else:
                    raise Exception("No valid images found")
                    
            except Exception as e:
                # Fallback or error
                logger.error(f"JPG to PDF failed: {e}")
                raise e

            # Serve File
            with open(output_pdf_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{final_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting files: {str(e)}")
            return redirect('jpg_to_pdf_tool')

    return render(request, 'core/jpg_to_pdf.html')

def sign_pdf_tool(request):
    """
    View to handle Free Sign PDF tool.
    Allows user to upload a PDF and a signature (or draw one), and overlays it.
    MVP: Places signature at bottom-right of the last page.
    """
    if request.method == 'POST':
        # PDF File
        try:
            pdf_file = request.FILES.get('pdf_file')
            signature_data = request.POST.get('signature_data') # Base64 string from canvas
            
            if not pdf_file or not signature_data:
                messages.error(request, "Please provide both a PDF and a signature.")
                return redirect('sign_pdf_tool')

            temp_files_to_clean = []
            
            # Save PDF
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in pdf_file.chunks():
                    tmp_pdf.write(chunk)
                temp_pdf_path = tmp_pdf.name
                temp_files_to_clean.append(temp_pdf_path)
            
            # Save Signature Image (Base64 -> PNG)
            import base64
            from PIL import Image
            import io
            
            format, imgstr = signature_data.split(';base64,') 
            ext = format.split('/')[-1] 
            
            signature_bytes = base64.b64decode(imgstr)
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_sig:
               tmp_sig.write(signature_bytes)
               temp_sig_path = tmp_sig.name
               temp_files_to_clean.append(temp_sig_path)

            # Output Path
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_pdf_path = tmp_out.name
                temp_files_to_clean.append(output_pdf_path)
            
            # Overlay Logic using PyMuPDF (fitz)
            doc = fitz.open(temp_pdf_path)
            
            # Target Page: Last Page
            page = doc[-1]
            rect = page.rect
            
            # Defined Size for Signature (e.g., 200px width, aspect ratio maintained)
            # Or simpler: Fixed box at bottom right
            # PDF coords: (0,0) is top-left in PyMuPDF? No, it varies. PyMuPDF is top-left usually.
            # Let's say we want it 50 units from bottom, 50 units from right.
            
            SIG_WIDTH = 150
            SIG_HEIGHT = 80 # Approx
            
            # Calculate position: Bottom Right
            x1 = rect.width - SIG_WIDTH - 50
            y1 = rect.height - SIG_HEIGHT - 50
            x2 = x1 + SIG_WIDTH
            y2 = y1 + SIG_HEIGHT
            
            signature_rect = fitz.Rect(x1, y1, x2, y2)
            
            # Insert Image
            page.insert_image(signature_rect, filename=temp_sig_path)
            
            doc.save(output_pdf_path)
            
            # Respond
            with open(output_pdf_path, 'rb') as f:
                pdf_data = f.read()
            
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{pdf_file.name.replace(".pdf", "")}_signed.pdf"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error signing PDF: {str(e)}")
            return redirect('sign_pdf_tool')

    return render(request, 'core/sign_pdf.html')

def html_to_pdf_tool(request):
    """
    View to handle Free HTML to PDF tool.
    Converts uploaded HTML files or URL to PDF.
    Note: This is a simplified implementation using reportlab.
    For complex HTML rendering, consider using a headless browser.
    """
    import requests
    from bs4 import BeautifulSoup
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    
    if request.method == 'POST':
        conversion_type = request.POST.get('conversion_type') # 'url' or 'file'
        
        temp_files_to_clean = []
        source_html = ""
        filename_prefix = "converted"
        
        try:
            if conversion_type == 'url':
                url = request.POST.get('url')
                if not url:
                    messages.error(request, "Please enter a valid URL.")
                    return redirect('html_to_pdf_tool')
                
                # Fetch URL content
                try:
                    headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
                    response = requests.get(url, headers=headers, timeout=30)
                    response.raise_for_status()
                    source_html = response.text
                    filename_prefix = url.split("//")[-1].replace("/", "_")[:20]
                except Exception as e:
                    messages.error(request, f"Failed to fetch URL: {str(e)}")
                    return redirect('html_to_pdf_tool')

            elif conversion_type == 'file':
                files = request.FILES.getlist('html_files')
                if not files:
                    messages.error(request, "Please upload an HTML file.")
                    return redirect('html_to_pdf_tool')
                
                uploaded_file = files[0]
                source_html = uploaded_file.read().decode('utf-8', errors='ignore')
                filename_prefix = uploaded_file.name.replace('.html', '').replace('.htm', '')
                
            else:
                 messages.error(request, "Invalid request.")
                 return redirect('html_to_pdf_tool')

            # Create PDF using reportlab (simple text extraction)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                output_path = tmp_pdf.name
                temp_files_to_clean.append(output_path)
            
            # Extract text from HTML
            soup = BeautifulSoup(source_html, 'html.parser')
            text_content = soup.get_text(separator='\n', strip=True)
            
            # Create PDF
            doc = SimpleDocTemplate(output_path, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Split into paragraphs and add to story
            for line in text_content.split('\n'):
                if line.strip():
                    try:
                        p = Paragraph(line, styles['BodyText'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                    except:
                        # Skip problematic lines
                        continue
                        
            doc.build(story)
                
            # Serve File
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
            
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename_prefix}.pdf"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response

        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error converting: {str(e)}")
            return redirect('html_to_pdf_tool')

    return render(request, 'core/html_to_pdf.html')

def rotate_pdf_tool(request):
    """
    View to handle Free Rotate PDF tool.
    Rotates all pages of a PDF by 90, 180, or 270 degrees.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files') # Support multiple if needed, but UI usually single
        rotation = request.POST.get('rotation') # 90, 180, 270 (clockwise)
        
        if not files or not rotation:
            messages.error(request, "Please provide a PDF and rotation angle.")
            return redirect('rotate_pdf_tool')

        try:
            rotation_angle = int(rotation)
            if rotation_angle not in [90, 180, 270]:
                raise ValueError("Invalid rotation")
        except:
             messages.error(request, "Invalid rotation angle.")
             return redirect('rotate_pdf_tool')
             
        temp_files_to_clean = []
        
        try:
            # Handle Single File for now (simplest for rotation)
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            doc = fitz.open(input_path)
            
            for page in doc:
                page.set_rotation(page.rotation + rotation_angle)
            
            filename_suffix = f"_rotated_{rotation_angle}"
            output_filename = uploaded_file.name.replace('.pdf', '') + filename_suffix + ".pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            doc.save(output_path)
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error rotating PDF: {str(e)}")
            return redirect('rotate_pdf_tool')
            
    return render(request, 'core/rotate_pdf.html')

def add_watermark_tool(request):
    """
    View to handle Free Add Watermark tool.
    Adds text watermark to all pages of a PDF.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        watermark_text = request.POST.get('watermark_text', 'CONFIDENTIAL')
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('add_watermark_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0] # Single file support for now
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            doc = fitz.open(input_path)
            
            for page in doc:
                # Calculate center
                rect = page.rect
                center = fitz.Point(rect.width / 2, rect.height / 2)
                
                # Insert Text
                # fontname="helv", fontsize=50, rotate=45, color=(0.5, 0.5, 0.5), opacity=0.3
                page.insert_text(
                    center,
                    watermark_text,
                    fontname="helv",
                    fontsize=50,
                    rotate=0,
                    color=(0.5, 0.5, 0.5)
                )
                


            output_filename = uploaded_file.name.replace('.pdf', '') + "_watermarked.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            doc.save(output_path)
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error creating watermark: {str(e)}")
            return redirect('add_watermark_tool')
            
    return render(request, 'core/add_watermark.html')

def protect_pdf_tool(request):
    """
    View to handle Free Protect PDF tool.
    Encrypts PDF with a password using pikepdf.
    """
    import pikepdf
    
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        password = request.POST.get('password')
        
        if not files or not password:
            messages.error(request, "Please provide a PDF and a password.")
            return redirect('protect_pdf_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            # Encrypt
            output_filename = uploaded_file.name.replace('.pdf', '') + "_protected.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            with pikepdf.Pdf.open(input_path) as pdf:
                pdf.save(output_path, encryption=pikepdf.Encryption(owner=password, user=password, R=6))
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error protecting PDF: {str(e)}")
            return redirect('protect_pdf_tool')
            
    return render(request, 'core/protect_pdf.html')

def unlock_pdf_tool(request):
    """
    View to handle Free Unlock PDF tool.
    Removes password from a PDF using pikepdf.
    """
    import pikepdf
    
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        password = request.POST.get('password', '')
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('unlock_pdf_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            output_filename = uploaded_file.name.replace('.pdf', '') + "_unlocked.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            try:
                # Try opening. If password needed and not provided, it will raise PasswordError
                with pikepdf.Pdf.open(input_path, password=password) as pdf:
                    pdf.save(output_path)
            except pikepdf.PasswordError:
                messages.error(request, "Incorrect password or password required.")
                return redirect('unlock_pdf_tool')
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error unlocking PDF: {str(e)}")
            return redirect('unlock_pdf_tool')
            
    return render(request, 'core/unlock_pdf.html')

def add_page_numbers_tool(request):
    """
    View to handle Free Add Page Numbers tool.
    Adds 'Page X of Y' to the bottom of all pages.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        position = request.POST.get('position', 'bottom-center') # simple for now
        
        if not files:
            messages.error(request, "Please upload a PDF file.")
            return redirect('add_page_numbers_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            for i, page in enumerate(doc):
                page_num = i + 1
                text = f"Page {page_num} of {total_pages}"
                
                rect = page.rect
                # Default: Bottom Center
                # y = rect.height - 20
                # x = rect.width / 2
                
                point = fitz.Point(rect.width / 2, rect.height - 20)
                
                # Textbox approach for better centering:
                footer_rect = fitz.Rect(0, rect.height - 40, rect.width, rect.height - 5)
                page.insert_textbox(footer_rect, text, fontsize=10, fontname="helv", align=1)

            output_filename = uploaded_file.name.replace('.pdf', '') + "_numbered.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            doc.save(output_path)
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error adding page numbers: {str(e)}")
            return redirect('add_page_numbers_tool')
            
    return render(request, 'core/add_page_numbers.html')

def remove_pages_tool(request):
    """
    View to handle Free Remove Pages tool.
    Deletes specified pages from a PDF.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        pages_to_remove_str = request.POST.get('pages_to_remove', '')
        
        if not files or not pages_to_remove_str:
            messages.error(request, "Please upload a PDF and specify pages to remove.")
            return redirect('remove_pages_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            # Parse Page Numbers
            # Input format: "1, 3-5, 7" -> remove these 1-based indexes
            pages_to_delete = set()
            try:
                parts = [p.strip() for p in pages_to_remove_str.split(',')]
                for part in parts:
                    if '-' in part:
                        start, end = part.split('-')
                        start, end = int(start), int(end)
                        for p in range(start, end + 1):
                            if 1 <= p <= total_pages:
                                pages_to_delete.add(p - 1) # 0-indexed
                    else:
                        p = int(part)
                        if 1 <= p <= total_pages:
                            pages_to_delete.add(p - 1)
            except ValueError:
                 messages.error(request, "Invalid page number format. Use '1, 3-5'.")
                 return redirect('remove_pages_tool')
            
            if len(pages_to_delete) == total_pages:
                 messages.error(request, "Cannot remove all pages.")
                 return redirect('remove_pages_tool')
            
            # Create new PDF with ONLY kept pages
            # Removing pages in fitz: doc.delete_pages(list)
            doc.delete_pages(list(pages_to_delete))
            
            output_filename = uploaded_file.name.replace('.pdf', '') + "_removed.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            doc.save(output_path)
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error removing pages: {str(e)}")
            return redirect('remove_pages_tool')
            
    return render(request, 'core/remove_pages.html')

def extract_pages_tool(request):
    """
    View to handle Free Extract Pages tool.
    Extracts specified pages from a PDF.
    """
    if request.method == 'POST':
        files = request.FILES.getlist('pdf_files')
        pages_to_extract_str = request.POST.get('pages_to_extract', '')
        
        if not files or not pages_to_extract_str:
            messages.error(request, "Please upload a PDF and specify pages to extract.")
            return redirect('extract_pages_tool')

        temp_files_to_clean = []
        
        try:
            uploaded_file = files[0]
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
                for chunk in uploaded_file.chunks():
                    tmp_pdf.write(chunk)
                input_path = tmp_pdf.name
                temp_files_to_clean.append(input_path)
            
            doc = fitz.open(input_path)
            total_pages = len(doc)
            
            # Parse Page Numbers
            pages_to_keep = [] # Ordered list for select
            try:
                parts = [p.strip() for p in pages_to_extract_str.split(',')]
                for part in parts:
                    if '-' in part:
                        start, end = part.split('-')
                        start, end = int(start), int(end)
                        for p in range(start, end + 1):
                            if 1 <= p <= total_pages:
                                pages_to_keep.append(p - 1)
                    else:
                        p = int(part)
                        if 1 <= p <= total_pages:
                            pages_to_keep.append(p - 1)
            except ValueError:
                 messages.error(request, "Invalid page number format. Use '1, 3-5'.")
                 return redirect('extract_pages_tool')
            
            if not pages_to_keep:
                 messages.error(request, "No valid pages selected.")
                 return redirect('extract_pages_tool')
            
            # Select pages (keep only these)
            doc.select(pages_to_keep)
            
            output_filename = uploaded_file.name.replace('.pdf', '') + "_extracted.pdf"
            
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_out:
                output_path = tmp_out.name
                temp_files_to_clean.append(output_path)
            
            doc.save(output_path)
            
            with open(output_path, 'rb') as f:
                pdf_data = f.read()
                
            response = HttpResponse(pdf_data, content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'
            
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            return response
            
        except Exception as e:
            for path in temp_files_to_clean:
                if os.path.exists(path):
                    os.remove(path)
            messages.error(request, f"Error extracting pages: {str(e)}")
            return redirect('extract_pages_tool')
            
    return render(request, 'core/extract_pages.html')

def whiteboard_tool(request):
    """
    View to handle the Whiteboard tool.
    Renders a full-screen or large canvas for drawing.
    """
    return render(request, 'core/whiteboard.html')
@login_required(login_url='order_panel_login')
def order_panel_send_notification(request):
    """
    Send notification to individual freelancer or broadcast to all
    """
    if request.method == 'POST':
        title = request.POST.get('title')
        message = request.POST.get('message')
        notification_type = request.POST.get('type', 'info')
        is_broadcast = request.POST.get('broadcast') == 'yes'
        
        if is_broadcast:
            # Send to all freelancers
            freelancers = Freelancer.objects.all()
            for freelancer in freelancers:
                FreelancerNotification.objects.create(
                    freelancer=freelancer,
                    title=title,
                    message=message,
                    notification_type=notification_type,
                    is_broadcast=True,
                    created_by=request.user
                )
            messages.success(request, f"Notification sent to {freelancers.count()} freelancers!")
        else:
            # Send to specific freelancer
            freelancer_id = request.POST.get('freelancer_id')
            if freelancer_id:
                freelancer = get_object_or_404(Freelancer, id=freelancer_id)
                FreelancerNotification.objects.create(
                    freelancer=freelancer,
                    title=title,
                    message=message,
                    notification_type=notification_type,
                    created_by=request.user
                )
                messages.success(request, f"Notification sent to {freelancer.name}!")
        
        # Redirect based on referer
        return redirect(request.META.get('HTTP_REFERER', 'order_panel_freelancers'))
    
    return redirect('order_panel_freelancers')
